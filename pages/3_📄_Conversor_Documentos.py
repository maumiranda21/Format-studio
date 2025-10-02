import streamlit as st
from pdf2docx import Converter
from tabula import read_pdf
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
import pypandoc
import io
import os
import pandas as pd

st.set_page_config(layout="wide", page_title="Conversor de Documentos")

st.title("📄 Herramienta 3: Conversor de Documentos")
st.markdown("Convierte entre formatos de documentos populares.")

tab1, tab2 = st.tabs(["🔄 PDF a Office", "✍️ Formatos de Texto (TXT, DOCX, MD...)"])

with tab1:
    st.header("Convertidor de PDF a Formatos de Office")
    uploaded_pdf = st.file_uploader("Sube tu archivo PDF", type="pdf", key="pdf_office")
    conversion_type = st.selectbox("Convertir a:",["PDF a Word (.docx)", "PDF a Excel (Tablas)", "PDF a PowerPoint (.pptx)"])
    if st.button("Convertir PDF", key="convert_pdf_btn"):
        if uploaded_pdf:
            with st.spinner("Realizando la conversión..."):
                pdf_bytes = uploaded_pdf.getvalue()
                pdf_name = os.path.splitext(uploaded_pdf.name)[0]
                if conversion_type == "PDF a Word (.docx)":
                    docx_buffer = io.BytesIO()
                    cv = Converter(stream=io.BytesIO(pdf_bytes))
                    cv.convert(docx_buffer)
                    cv.close()
                    st.success("¡PDF convertido a Word!")
                    st.download_button(label="📥 Descargar DOCX", data=docx_buffer.getvalue(), file_name=f"{pdf_name}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                elif conversion_type == "PDF a Excel (Tablas)":
                    tables = read_pdf(io.BytesIO(pdf_bytes), pages='all', multiple_tables=True)
                    if not tables: st.warning("No se encontraron tablas en el PDF.")
                    else:
                        excel_buffer = io.BytesIO()
                        with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
                            for i, df in enumerate(tables): df.to_excel(writer, sheet_name=f'Tabla_{i+1}', index=False)
                        st.success(f"¡Se extrajeron {len(tables)} tablas a Excel!")
                        st.download_button(label="📥 Descargar XLSX", data=excel_buffer.getvalue(), file_name=f"{pdf_name}_tablas.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                elif conversion_type == "PDF a PowerPoint (.pptx)":
                    images = convert_from_bytes(pdf_bytes)
                    prs = Presentation()
                    for image in images:
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        img_buffer = io.BytesIO()
                        image.save(img_buffer, format='PNG')
                        slide.shapes.add_picture(io.BytesIO(img_buffer.getvalue()), Inches(0.5), Inches(0.75), width=Inches(9))
                    pptx_buffer = io.BytesIO()
                    prs.save(pptx_buffer)
                    st.success("¡PDF convertido a PowerPoint!")
                    st.download_button(label="📥 Descargar PPTX", data=pptx_buffer.getvalue(), file_name=f"{pdf_name}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        else: st.warning("Por favor, sube un archivo PDF.")

with tab2:
    st.header("Convertidor entre Formatos de Texto")
    uploaded_text_file = st.file_uploader("Sube tu archivo (TXT, DOCX, MD, RTF)", type=["txt", "docx", "md", "rtf"], key="text_formats")
    col1, col2 = st.columns(2)
    with col1: from_format = st.selectbox("Formato de Origen", ["docx", "rtf", "md", "txt"])
    with col2: to_format = st.selectbox("Formato de Destino", ["docx", "rtf", "md", "txt", "pdf"])
    if st.button("Convertir Texto", key="convert_text_btn"):
        if uploaded_text_file:
            with st.spinner("Convirtiendo..."):
                file_content = uploaded_text_file.read().decode("utf-8")
                output = pypandoc.convert_text(file_content, to_format, format=from_format)
                output_filename = f"{os.path.splitext(uploaded_text_file.name)[0]}.{to_format}"
                st.success("¡Conversión completada!")
                st.download_button(label=f"📥 Descargar {output_filename}", data=output.encode('utf-8') if isinstance(output, str) else output, file_name=output_filename)
        else: st.warning("Por favor, sube un archivo.")
